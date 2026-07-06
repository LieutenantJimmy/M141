#!/usr/bin/env python3
"""
M141 LB3 - Praesentations-Deck-Generator (reproduzierbar)
==========================================================
Erzeugt LB3_Praesentation.pptx (8 Folien, 16:9, praesentierenden-neutral).

Nutzung:
    pip install python-pptx
    python3 build_pptx.py            # schreibt ./LB3_Praesentation.pptx

Das Deck ist bewusst schlank: visuelles Geruest fuer die Live-Demo,
kein Ersatz fuer die Doku (siehe README_Praxisarbeit.md, VERIFICATION.md).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ----- Farbwelt (dunkles Petrol + neutrale Grautoene) ------------------------
INK    = RGBColor(0x1B, 0x2A, 0x33)   # fast-schwarz
TEAL   = RGBColor(0x0F, 0x76, 0x6E)   # Akzent
LIGHT  = RGBColor(0xF4, 0xF6, 0xF7)   # Folienhintergrund content
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x5A, 0x6A, 0x72)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
RED    = RGBColor(0xB0, 0x3A, 0x2E)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame

def para(tf, text, size, color=INK, bold=False, first=False, align=PP_ALIGN.LEFT,
         bullet=None, space=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    run = p.add_run()
    run.text = (f"{bullet}  {text}" if bullet else text)
    f = run.font
    f.size, f.bold, f.color.rgb, f.name = Pt(size), bold, color, "Calibri"
    return p

def accent_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.16))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()

def content_slide(title, kicker=None):
    s = prs.slides.add_slide(BLANK)
    bg(s, LIGHT); accent_bar(s)
    tf = box(s, Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.1))
    if kicker:
        para(tf, kicker, 13, TEAL, bold=True, first=True, space=2)
        para(tf, title, 30, INK, bold=True)
    else:
        para(tf, title, 30, INK, bold=True, first=True)
    return s

def footer(s, text="M141 LB3 · Backpacker_LB3 · eigene Cloud-DB"):
    tf = box(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35))
    para(tf, text, 10, GREY, first=True)

# ============================ Folie 1 - Titel ================================
s = prs.slides.add_slide(BLANK); bg(s, INK)
bar = s.shapes.add_shape(1, Inches(0), Inches(4.55), W, Inches(0.10))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
tf = box(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.2))
para(tf, "M141 · LB3 Praxisarbeit", 16, TEAL, bold=True, first=True, space=10)
para(tf, "Backpacker_LB3 — Migration in die eigene Cloud", 40, WHITE, bold=True, space=12)
para(tf, "Von Access zu MariaDB: lokal aufgebaut, gehärtet und in eine "
         "selbst gehostete Homelab-Cloud migriert — statt Managed-Provider.",
     18, RGBColor(0xC9, 0xD4, 0xD9))
tf2 = box(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5))
para(tf2, "TBZ · Modul 141 «DB-Systeme in Betrieb nehmen» · Demo Juli 2026", 12,
     RGBColor(0x8F, 0xA3, 0xAB), first=True)

# ============================ Folie 2 - Architektur ==========================
s = content_slide("Architektur", "SETUP")
tf = box(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.0))
para(tf, "CSV (Access)   ➜   lokale MariaDB (Staging → Bereinigung → 2.NF)   ➜   "
         "mysqldump   ➜  TLS  ➜   eigene Cloud-DB", 17, TEAL, bold=True,
     first=True, align=PP_ALIGN.CENTER)
tf = box(s, Inches(1.0), Inches(2.7), Inches(11.2), Inches(3.8))
for t in [
    "Cloud = LXC-Container «cloud-db-giovanni» (CT 9003) auf dem eigenen Proxmox-Host",
    "Endpoint 192.168.1.62:3306 · MariaDB 11.8 · utf8mb4 · InnoDB · 5 Foreign Keys",
    "TLS erzwungen (require_secure_transport=ON, TLSv1.3, eigene CA mit SAN)",
    "IP-Allowlist per Proxmox-Firewall — nur berechtigte Quell-IPs, sonst DROP",
    "Zugriff ausschliesslich per mysql-CLI mit CA-Verifikation (--ssl-ca)",
]:
    para(tf, t, 18, INK, bullet="▸", space=12, first=(t.startswith("Cloud")))
footer(s)

# ============================ Folie 3 - DCL-Modell ===========================
s = content_slide("Zugriffskontrolle (DCL)", "WER DARF WAS")
tf = box(s, Inches(1.0), Inches(1.7), Inches(11.2), Inches(4.6))
para(tf, "2 Rollen + 3 User — alle Cloud-User mit REQUIRE SSL:", 18, INK,
     bold=True, first=True, space=12)
for t, c in [
    ("role_benutzer (Empfang):  Buchungen/Positionen voll · Gäste ohne DELETE · "
     "Spalte «Password» NIE lesbar · «deaktiviert» nur lesbar", INK),
    ("role_management:  alles voll — ausser Buchungen/Positionen = nur SELECT", INK),
    ("giovanni_dba:  ALL PRIVILEGES auf die DB (Administration)", INK),
    ("→ Zugriffsmatrix bis auf SPALTEN-Ebene erzwungen — Verstösse enden in "
     "ERROR 1142 / 1143 (live vorführbar)", TEAL),
]:
    para(tf, t, 17, c, bullet="▸", space=14)
footer(s)

# ============================ Folie 4 - Migration ============================
s = content_slide("Automatisierte Migration (DDL + DML)", "LOKAL ➜ CLOUD")
tf = box(s, Inches(1.0), Inches(1.7), Inches(11.2), Inches(3.2))
for i, t in enumerate([
    "Ein Skript: DB anlegen → Dump einspielen (Struktur + Daten, --single-transaction) "
    "→ DCL anwenden → Smoke-Test — alles über TLS, < 1 Minute",
    "Idempotent: kompletter Lauf ZWEIMAL ausgeführt — identisches Ergebnis, 0 Fehler "
    "(Beweis: VERIFICATION.md §1)",
    "Personalisierter Migrations-Testdatensatz «Giovanni-Test» in der Cloud verifiziert",
]):
    para(tf, t, 18, INK, bullet="▸", space=14, first=(i == 0))
tf = box(s, Inches(1.0), Inches(5.0), Inches(11.2), Inches(1.4))
para(tf, "Row-Counts lokal = Cloud:", 16, GREY, bold=True, first=True, space=4)
para(tf, "personen 2036 · benutzer 11 · land 82 · leistung 8 · buchung 1006 · positionen 1746",
     19, GREEN, bold=True)
footer(s)

# ============================ Folie 5 - Haertung =============================
s = content_slide("Härtung — 8 Punkte", "SICHERHEIT")
left = ["TLS erzwungen + REQUIRE SSL je User",
        "Eigene CA + Server-Zertifikat (SAN auf Endpoint-IP)",
        "IP-Allowlist — kein 0.0.0.0/0",
        "Default-Deny (policy_in: DROP)"]
right = ["local-infile = 0 (kein LOAD DATA LOCAL)",
         "skip-name-resolve (keine DNS-Lookups)",
         "slow_query_log + error_log aktiv",
         "Least Privilege: Admin- vs. App-User getrennt"]
tf = box(s, Inches(0.9), Inches(1.9), Inches(5.9), Inches(4.4))
for i, t in enumerate(left):
    para(tf, t, 18, INK, bullet=f"{i+1}.", space=18, first=(i == 0))
tf = box(s, Inches(6.9), Inches(1.9), Inches(5.9), Inches(4.4))
for i, t in enumerate(right):
    para(tf, t, 18, INK, bullet=f"{i+5}.", space=18, first=(i == 0))
tf = box(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.6))
para(tf, "Jeder Punkt einzeln live belegt — VERIFICATION.md §3", 14, TEAL,
     bold=True, first=True)
footer(s)

# ============================ Folie 6 - Tests ================================
s = content_slide("Testergebnisse — alles grün", "QUALITÄT")
tf = box(s, Inches(1.0), Inches(1.7), Inches(11.2), Inches(4.6))
for i, (t, c) in enumerate([
    ("Lokal: 13 Konsistenz-Checks (FKs, Sentinels, Duplikate) + 19 Rollen-Tests "
     "positiv/negativ — alle ✓", GREEN),
    ("Cloud: Counts identisch · 5 FKs · utf8mb4 · Rollen greifen · "
     "Klartext abgewiesen (ERROR 3159)", GREEN),
    ("Negativ-Tests treffen exakt ein: 1142 (Tabelle verweigert) · "
     "1143 (Spalte verweigert) · 3159 (TLS-Pflicht)", GREEN),
    ("Ehrlich erklärt: 22 Positionen mit negativer Anzahl/Preis = BEWUSST erlaubte "
     "Stornos (dokumentiert) — echte CHECK-Verletzungen: 0", INK),
    ("Jeder Nachweis doppelt: Screenshot + rohe Konsolen-Ausgabe im Repo", GREY),
]):
    para(tf, t, 17, c, bullet="▸", space=14, first=(i == 0))
footer(s)

# ============================ Folie 7 - Warum eigene Cloud ===================
s = content_slide("Warum eigene Cloud statt Aiven/AWS?", "ENTSCHEID")
tf = box(s, Inches(1.0), Inches(1.7), Inches(11.2), Inches(4.6))
for i, (t, c) in enumerate([
    ("Rahmen: «Andere oder eigene Cloud-DB gibt +» — die eigene Cloud ist die "
     "Max-Bonus-Option", TEAL),
    ("AWS: kein TBZ-Schulungsaccount · Aiven: evaluiert (gewichtete Matrix, "
     "MS_A_Cloud_Evaluation.md)", INK),
    ("Eigene Cloud zeigt, was Managed-Dienste verstecken: TLS-Setup, Zertifikate, "
     "Firewall, Härtung — alles selbst gebaut, dokumentiert, reproduzierbar", INK),
    ("Trade-off transparent: Betrieb in Eigenverantwortung — inkl. real erlebtem "
     "Host-Ausfall und host-agnostischem Re-Deployment", GREY),
]):
    para(tf, t, 18, c, bullet="▸", space=16, first=(i == 0))
footer(s)

# ============================ Folie 8 - Live-Demo ============================
s = content_slide("Live-Demo", "JETZT")
tf = box(s, Inches(1.0), Inches(1.7), Inches(11.2), Inches(3.9))
for i, t in enumerate([
    "Preflight:  ein Befehl, 21 Checks, Auto-Heal  →  GO / NO-GO",
    "3 User verbinden per TLS:  benutzer · manager · dba",
    "Pro User: 1 erlaubte + 1 verbotene Aktion (ERROR 1142/1143 live)",
    "Finale: Klartext-Login scheitert — ERROR 3159 (TLS-Pflicht bewiesen)",
]):
    para(tf, t, 20, INK, bullet=f"{i+1}.", space=18, first=(i == 0))
tf = box(s, Inches(1.0), Inches(5.7), Inches(11.2), Inches(1.0))
para(tf, "Sicherheitsnetz: Golden Snapshot — Rollback in < 1 Minute.", 15, GREY,
     bold=True, first=True, space=4)
para(tf, "Fragen?", 26, TEAL, bold=True)
footer(s)

prs.save("LB3_Praesentation.pptx")
print(f"OK: LB3_Praesentation.pptx geschrieben ({len(prs.slides)} Folien)")
